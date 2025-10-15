# services/file.py
import os
import time
import mimetypes
from pathlib import Path
from typing import Optional

from fastapi import UploadFile
from loguru import logger
from PyPDF2 import PdfReader
import docx2txt
import csv
import pptx  # python-pptx

from models.models import Document, DocumentMetadata, Source

# Dossier persistant (monte-le dans docker-compose : volumes: - uploads:/data/uploads)
UPLOAD_DIR = Path("/data/uploads")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)


def _safe_filename(name: str) -> str:
    # petite sanitation sans dépendre de werkzeug
    name = name or "upload"
    name = name.replace("\\", "/").split("/")[-1]
    keep = "._-() "
    sanitized = "".join(ch for ch in name if ch.isalnum() or ch in keep)
    return sanitized or f"file_{int(time.time())}"


def _guess_mimetype(path: str, provided: Optional[str]) -> str:
    if provided:
        return provided
    mt, _ = mimetypes.guess_type(path)
    if mt:
        return mt
    if path.endswith(".md"):
        return "text/markdown"
    raise ValueError("Unsupported file type (cannot guess mimetype)")


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Retourne le texte du fichier à partir de son chemin."""
    mt = _guess_mimetype(filepath, mimetype)
    logger.info(f"extract_text_from_filepath: {filepath} ({mt})")

    if mt == "application/pdf":
        reader = PdfReader(filepath)
        texts = []
        for page in reader.pages:
            try:
                texts.append(page.extract_text() or "")
            except Exception:
                pass
        return " ".join(texts)

    if mt in ("text/plain", "text/markdown"):
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    if mt == "text/csv":
        out = []
        with open(filepath, "r", encoding="utf-8", errors="ignore", newline="") as f:
            reader = csv.reader(f)
            for row in reader:
                out.append(" ".join(row))
        return "\n".join(out)

    if mt == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return docx2txt.process(filepath) or ""

    if mt == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = []
        prs = pptx.Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text.append(" ".join(run.text or "" for run in para.runs))
        return "\n".join(text)

    raise ValueError(f"Unsupported file type: {mt}")


async def get_document_from_file(file: UploadFile, metadata: DocumentMetadata) -> Document:
    """Sauvegarde le fichier, renseigne metadata (url/source_id), puis extrait le texte."""

    # 1) Sauvegarde le flux reçu une seule fois
    orig_name = _safe_filename(file.filename or "upload")
    stamp = int(time.time())
    dst = UPLOAD_DIR / f"{stamp}_{orig_name}"

    blob = await file.read()
    with open(dst, "wb") as out:
        out.write(blob)

    # 2) Métadonnées (utiliser source_id au lieu de document_id)
    if metadata is None:
        metadata = DocumentMetadata(source=Source.file)
    if getattr(metadata, "source", None) is None:
        metadata.source = Source.file
    if not getattr(metadata, "url", None):
        metadata.url = f"file://{dst}"
    if not getattr(metadata, "source_id", None):
        metadata.source_id = dst.name  # <<— identifiant de fichier

    # 3) Extraction texte depuis le fichier sauvegardé
    try:
        extracted_text = extract_text_from_filepath(str(dst), file.content_type)
    except Exception as e:
        logger.error(f"extract_text failed for {dst}: {e}")
        raise

    return Document(text=extracted_text, metadata=metadata)
